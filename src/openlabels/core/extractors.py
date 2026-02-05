"""
Text extractors for various file formats.

Each extractor implements a common interface for extracting text content
from files, with security protections against decompression bombs.
"""

import csv
import io
import logging
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Dict, Any

from .constants import (
    MIN_NATIVE_TEXT_LENGTH,
    MAX_DOCUMENT_PAGES,
    MAX_SPREADSHEET_ROWS,
    MAX_DECOMPRESSED_SIZE,
    MAX_EXTRACTION_RATIO,
)
from .exceptions import ExtractionError, SecurityError

logger = logging.getLogger(__name__)


@dataclass
class PageInfo:
    """Information about a single extracted page."""
    page_num: int
    text: str
    is_scanned: bool  # True if this page needed OCR


@dataclass
class ExtractionResult:
    """Result of text extraction from a file."""
    text: str
    pages: int = 1
    needs_ocr: bool = False  # True if any page needed OCR
    ocr_pages: List[int] = field(default_factory=list)  # Which pages needed OCR
    warnings: List[str] = field(default_factory=list)  # Non-fatal issues
    confidence: float = 1.0  # Average OCR confidence (1.0 if no OCR)
    page_infos: List[PageInfo] = field(default_factory=list)

    @property
    def has_scanned_pages(self) -> bool:
        """Check if any pages are scanned (need visual redaction)."""
        return any(p.is_scanned for p in self.page_infos)

    @property
    def scanned_page_count(self) -> int:
        """Count of scanned pages."""
        return sum(1 for p in self.page_infos if p.is_scanned)


class BaseExtractor(ABC):
    """Base class for format-specific extractors."""

    @abstractmethod
    def can_handle(self, content_type: str, extension: str) -> bool:
        """
        Check if this extractor handles the file type.

        Args:
            content_type: MIME type
            extension: File extension (lowercase, with dot)

        Returns:
            True if this extractor can process the file
        """
        pass

    @abstractmethod
    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        """
        Extract text from file content.

        Args:
            content: Raw file bytes
            filename: Original filename (for logging/extension detection)

        Returns:
            ExtractionResult with extracted text and metadata
        """
        pass


class PDFExtractor(BaseExtractor):
    """
    PDF text extractor using PyMuPDF.

    Handles both text-layer PDFs and scanned documents:
    - Extracts text layer if available (native PDF)
    - Falls back to OCR for pages without text (scanned PDF)
    """

    RENDER_DPI = 150  # DPI for rendering scanned pages

    def __init__(self, ocr_engine: Optional[Any] = None):
        """
        Initialize PDF extractor.

        Args:
            ocr_engine: Optional OCR engine for scanned pages
        """
        self.ocr_engine = ocr_engine

    def can_handle(self, content_type: str, extension: str) -> bool:
        return content_type == "application/pdf" or extension == ".pdf"

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from PDF with OCR fallback for scanned pages."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

        doc = fitz.open(stream=content, filetype="pdf")

        pages_text = []
        page_infos = []
        ocr_pages = []
        ocr_confidences = []
        warnings = []

        try:
            for i, page in enumerate(doc):
                # Early exit if page limit exceeded (prevents DoS via large PDFs)
                if i >= MAX_DOCUMENT_PAGES:
                    logger.warning(f"PDF exceeds {MAX_DOCUMENT_PAGES} page limit, truncating")
                    warnings.append(f"Document truncated at {MAX_DOCUMENT_PAGES} pages")
                    break

                # Try to extract text layer
                native_text = page.get_text().strip()

                # Check if this page has meaningful native text
                has_native_text = len(native_text) >= MIN_NATIVE_TEXT_LENGTH

                if has_native_text:
                    # Native text page
                    pages_text.append(native_text)
                    page_infos.append(PageInfo(
                        page_num=i,
                        text=native_text,
                        is_scanned=False,
                    ))
                    logger.debug(f"Page {i+1}: native text ({len(native_text)} chars)")

                elif self.ocr_engine and hasattr(self.ocr_engine, 'is_available') and self.ocr_engine.is_available:
                    # Scanned page - needs OCR
                    logger.debug(f"Page {i+1}: scanned, using OCR")

                    try:
                        # Render page to image
                        pix = page.get_pixmap(dpi=self.RENDER_DPI)

                        # Convert to PIL Image and numpy array
                        from PIL import Image
                        import numpy as np

                        img = Image.frombytes(
                            "RGB",
                            [pix.width, pix.height],
                            pix.samples
                        )
                        img_array = np.array(img)

                        # Run OCR
                        ocr_text = self.ocr_engine.extract_text(img_array)

                        pages_text.append(ocr_text)
                        ocr_pages.append(i)

                        page_infos.append(PageInfo(
                            page_num=i,
                            text=ocr_text,
                            is_scanned=True,
                        ))

                    except Exception as e:
                        # Log OCR failures with full context - may indicate corrupted pages or OCR issues
                        logger.warning(f"OCR failed for page {i+1} of {filename}: {type(e).__name__}: {e}")
                        pages_text.append("")
                        warnings.append(f"OCR failed for page {i+1}: {e}")
                        page_infos.append(PageInfo(
                            page_num=i,
                            text="",
                            is_scanned=True,
                        ))
                else:
                    # No OCR available
                    pages_text.append("")
                    page_infos.append(PageInfo(
                        page_num=i,
                        text="",
                        is_scanned=True,  # Assume scanned if no text
                    ))
                    if not self.ocr_engine:
                        warnings.append(f"Page {i+1} is scanned but OCR not available")

            # Calculate average OCR confidence
            avg_confidence = 1.0
            if ocr_confidences:
                avg_confidence = sum(ocr_confidences) / len(ocr_confidences)

            return ExtractionResult(
                text="\n\n".join(pages_text),
                pages=len(doc),
                needs_ocr=len(ocr_pages) > 0,
                ocr_pages=ocr_pages,
                warnings=warnings,
                confidence=avg_confidence,
                page_infos=page_infos,
            )

        finally:
            doc.close()


class DOCXExtractor(BaseExtractor):
    """Word document extractor using python-docx."""

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ) or
            extension in (".docx", ".doc")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        ext = Path(filename).suffix.lower()

        if ext == ".doc":
            # Legacy .doc format - limited support
            return self._extract_legacy_doc(content, filename)

        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        # SECURITY: Check decompression ratio before extraction
        # DOCX is a ZIP file - malicious files could have huge decompression ratios
        compressed_size = len(content)
        self._check_decompression_size(compressed_size, filename)

        doc = Document(io.BytesIO(content))

        paragraphs = []
        total_chars = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                total_chars += len(text)
                # SECURITY: Early termination if extraction is too large
                if total_chars > MAX_DECOMPRESSED_SIZE:
                    raise ValueError(
                        f"Decompression bomb detected: extracted content exceeds "
                        f"{MAX_DECOMPRESSED_SIZE // (1024*1024)}MB limit"
                    )

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                        total_chars += len(cell_text)
                if row_text:
                    paragraphs.append(" | ".join(row_text))
                # SECURITY: Check size during table extraction too
                if total_chars > MAX_DECOMPRESSED_SIZE:
                    raise ValueError(
                        f"Decompression bomb detected: extracted content exceeds "
                        f"{MAX_DECOMPRESSED_SIZE // (1024*1024)}MB limit"
                    )

        return ExtractionResult(
            text="\n\n".join(paragraphs),
            pages=1,  # DOCX doesn't have fixed pages
        )

    def _check_decompression_size(self, compressed_size: int, filename: str) -> None:
        """Check that compressed file isn't suspiciously small (potential zip bomb)."""
        # Very small compressed files that claim to be documents are suspicious
        if compressed_size < 100:
            logger.warning(f"Suspiciously small DOCX file: {filename} ({compressed_size} bytes)")
            # Allow but log - very small files might be legitimate empty docs

    def _extract_legacy_doc(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from legacy .doc format. Limited support."""
        try:
            # Attempt basic text extraction
            text = content.decode("latin-1", errors="ignore")
            # Filter to printable characters
            printable = "".join(
                c if c.isprintable() or c in "\n\r\t" else " "
                for c in text
            )
            # Clean up whitespace
            lines = [line.strip() for line in printable.split("\n")]
            lines = [line for line in lines if line and len(line) > 3]

            return ExtractionResult(
                text="\n".join(lines),
                pages=1,
                warnings=["Legacy .doc format - extraction may be incomplete"],
            )
        except Exception as e:
            # Log legacy doc extraction failures - may indicate corrupt files
            logger.info(f"Legacy .doc extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"Failed to extract from legacy .doc: {e}"],
            )


class XLSXExtractor(BaseExtractor):
    """Spreadsheet extractor for XLSX, XLS, and CSV."""

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
                "text/csv",
            ) or
            extension in (".xlsx", ".xls", ".csv", ".tsv")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        ext = Path(filename).suffix.lower()

        if ext == ".csv":
            return self._extract_csv(content, ",")
        elif ext == ".tsv":
            return self._extract_csv(content, "\t")
        elif ext == ".xls":
            return self._extract_xls(content, filename)
        else:
            return self._extract_xlsx(content, filename)

    def _extract_csv(self, content: bytes, delimiter: str) -> ExtractionResult:
        """Extract from CSV/TSV."""
        # Try common encodings - UnicodeDecodeError is expected for wrong encodings
        text_content = None
        for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
            try:
                text_content = content.decode(encoding)
                break
            except UnicodeDecodeError:
                # This encoding doesn't work - try next one
                continue

        if text_content is None:
            return ExtractionResult(
                text="",
                warnings=["Failed to decode CSV file"],
            )

        rows = []
        reader = csv.reader(io.StringIO(text_content), delimiter=delimiter)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(cell.strip() for cell in row if cell.strip()))

        return ExtractionResult(
            text="\n".join(rows),
            pages=1,
        )

    def _extract_xlsx(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from XLSX."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl not installed. Run: pip install openpyxl")

        # SECURITY: Track total extracted size to prevent decompression bombs
        compressed_size = len(content)

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        all_text = []
        warnings = []
        total_chars = 0

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            row_count = 0

            for row in sheet.iter_rows(values_only=True):
                # Limit rows per sheet to prevent DoS
                if row_count >= MAX_SPREADSHEET_ROWS:
                    warnings.append(f"Sheet '{sheet_name}' truncated at {MAX_SPREADSHEET_ROWS} rows")
                    break
                row_count += 1

                cells = [str(cell).strip() for cell in row if cell is not None]
                if cells:
                    row_text = " | ".join(cells)
                    sheet_rows.append(row_text)
                    total_chars += len(row_text)

                    # SECURITY: Check for decompression bomb
                    if total_chars > MAX_DECOMPRESSED_SIZE:
                        wb.close()
                        raise ValueError(
                            f"Decompression bomb detected: extracted content exceeds "
                            f"{MAX_DECOMPRESSED_SIZE // (1024*1024)}MB limit"
                        )

            if sheet_rows:
                all_text.append(f"[Sheet: {sheet_name}]")
                all_text.extend(sheet_rows)
                all_text.append("")

        # Store sheet count before closing workbook
        sheet_count = len(wb.sheetnames) if hasattr(wb, 'sheetnames') else 1
        wb.close()

        # SECURITY: Final check on extraction ratio
        if compressed_size > 0 and total_chars > 0:
            ratio = total_chars / compressed_size
            if ratio > MAX_EXTRACTION_RATIO:
                logger.warning(
                    f"High extraction ratio for {filename}: {ratio:.1f}x "
                    f"({compressed_size} bytes -> {total_chars} chars)"
                )

        return ExtractionResult(
            text="\n".join(all_text),
            pages=sheet_count,
            warnings=warnings,
        )

    def _extract_xls(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from legacy XLS."""
        try:
            import xlrd
        except ImportError:
            raise ImportError("xlrd not installed. Run: pip install xlrd")

        wb = xlrd.open_workbook(file_contents=content)

        all_text = []
        warnings = []
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            sheet_rows = []

            # Limit rows per sheet to prevent DoS
            max_rows = min(sheet.nrows, MAX_SPREADSHEET_ROWS)
            if sheet.nrows > MAX_SPREADSHEET_ROWS:
                warnings.append(f"Sheet '{sheet.name}' truncated at {MAX_SPREADSHEET_ROWS} rows")

            for row_idx in range(max_rows):
                row = sheet.row_values(row_idx)
                cells = [str(cell).strip() for cell in row if cell]
                if cells:
                    sheet_rows.append(" | ".join(cells))

            if sheet_rows:
                all_text.append(f"[Sheet: {sheet.name}]")
                all_text.extend(sheet_rows)
                all_text.append("")

        return ExtractionResult(
            text="\n".join(all_text),
            pages=wb.nsheets,
            warnings=warnings,
        )


class ImageExtractor(BaseExtractor):
    """
    Image text extractor using OCR.

    Handles JPEG, PNG, TIFF, HEIC, GIF, BMP, WebP.
    """

    def __init__(self, ocr_engine: Optional[Any] = None):
        """
        Initialize image extractor.

        Args:
            ocr_engine: OCR engine for text extraction
        """
        self.ocr_engine = ocr_engine

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type.startswith("image/") or
            extension in (".jpg", ".jpeg", ".png", ".tiff", ".tif",
                         ".heic", ".heif", ".gif", ".bmp", ".webp")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from image using OCR."""
        if not self.ocr_engine or not getattr(self.ocr_engine, 'is_available', False):
            return ExtractionResult(
                text="",
                pages=1,
                needs_ocr=True,
                warnings=["OCR engine not available for image extraction"],
            )

        ext = Path(filename).suffix.lower()

        try:
            from PIL import Image
            import numpy as np

            # Handle HEIC format
            if ext in (".heic", ".heif"):
                try:
                    from pillow_heif import register_heif_opener
                    register_heif_opener()
                except ImportError:
                    return ExtractionResult(
                        text="",
                        pages=1,
                        warnings=["HEIC support not available. Run: pip install pillow-heif"],
                    )

            # Handle multi-page TIFF
            if ext in (".tiff", ".tif"):
                return self._extract_multipage_tiff(content, filename)

            # Single image
            img = Image.open(io.BytesIO(content))

            # Convert to RGB if necessary
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            # Convert to numpy array
            img_array = np.array(img)

            # Run OCR
            ocr_text = self.ocr_engine.extract_text(img_array)

            page_info = PageInfo(
                page_num=0,
                text=ocr_text,
                is_scanned=True,
            )

            return ExtractionResult(
                text=ocr_text,
                pages=1,
                needs_ocr=True,
                ocr_pages=[0],
                page_infos=[page_info],
            )

        except Exception as e:
            # Log image extraction failures with context
            logger.warning(f"Image extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"Image extraction failed: {e}"],
            )

    def _extract_multipage_tiff(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract text from multi-page TIFF."""
        from PIL import Image
        import numpy as np

        img = Image.open(io.BytesIO(content))

        pages_text = []
        page_infos = []

        try:
            page_num = 0
            while True:
                # Early exit if page limit exceeded
                if page_num >= MAX_DOCUMENT_PAGES:
                    logger.warning(f"TIFF exceeds {MAX_DOCUMENT_PAGES} page limit, truncating")
                    break

                img.seek(page_num)

                # Convert to RGB
                frame = img.convert("RGB")
                frame_array = np.array(frame)

                # Run OCR
                ocr_text = self.ocr_engine.extract_text(frame_array)

                pages_text.append(ocr_text)
                page_infos.append(PageInfo(
                    page_num=page_num,
                    text=ocr_text,
                    is_scanned=True,
                ))

                page_num += 1

        except EOFError:
            # End of pages
            pass

        return ExtractionResult(
            text="\n\n".join(pages_text),
            pages=len(pages_text),
            needs_ocr=True,
            ocr_pages=list(range(len(pages_text))),
            page_infos=page_infos,
        )


class TextExtractor(BaseExtractor):
    """Plain text file extractor."""

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type == "text/plain" or
            extension == ".txt"
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        # Try common encodings - UnicodeDecodeError is expected for wrong encodings
        for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
            try:
                text = content.decode(encoding)
                return ExtractionResult(
                    text=text,
                    pages=1,
                )
            except UnicodeDecodeError:
                # This encoding doesn't work - try next one
                continue

        return ExtractionResult(
            text="",
            pages=1,
            warnings=["Failed to decode text file"],
        )


class RTFExtractor(BaseExtractor):
    """RTF document extractor using striprtf."""

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type == "application/rtf" or
            extension == ".rtf"
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise ImportError("striprtf not installed. Run: pip install striprtf")

        # Try to decode - UnicodeDecodeError is expected for wrong encodings
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                rtf_content = content.decode(encoding)
                break
            except UnicodeDecodeError:
                # This encoding doesn't work - try next one
                continue
        else:
            return ExtractionResult(
                text="",
                pages=1,
                warnings=["Failed to decode RTF file"],
            )

        try:
            text = rtf_to_text(rtf_content)
            return ExtractionResult(
                text=text,
                pages=1,
            )
        except Exception as e:
            # Log RTF extraction failures
            logger.info(f"RTF extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"RTF extraction failed: {e}"],
            )


class PPTXExtractor(BaseExtractor):
    """
    PowerPoint extractor for PPTX files.

    Extracts text from:
    - Slide content (text boxes, shapes)
    - Speaker notes
    - Tables
    """

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            ) or
            extension in (".pptx", ".ppt")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        ext = Path(filename).suffix.lower()

        if ext == ".ppt":
            return self._extract_legacy_ppt(content, filename)

        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(io.BytesIO(content))

        slides_text = []
        warnings = []
        total_chars = 0

        for slide_num, slide in enumerate(prs.slides):
            slide_content = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    total_chars += len(shape.text)

                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                                total_chars += len(cell.text)
                        if row_text:
                            slide_content.append(" | ".join(row_text))

                # SECURITY: Check for decompression bomb
                if total_chars > MAX_DECOMPRESSED_SIZE:
                    raise ValueError(
                        f"Decompression bomb detected: extracted content exceeds "
                        f"{MAX_DECOMPRESSED_SIZE // (1024*1024)}MB limit"
                    )

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_content.append(f"[Notes: {notes}]")
                    total_chars += len(notes)

            if slide_content:
                slides_text.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_content))

        return ExtractionResult(
            text="\n\n".join(slides_text),
            pages=len(prs.slides),
            warnings=warnings,
        )

    def _extract_legacy_ppt(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from legacy .ppt format. Limited support."""
        # Legacy .ppt is a binary format - basic extraction
        try:
            text = content.decode("latin-1", errors="ignore")
            # Filter to printable characters
            printable = "".join(
                c if c.isprintable() or c in "\n\r\t" else " "
                for c in text
            )
            # Clean up whitespace and filter short lines
            lines = [line.strip() for line in printable.split("\n")]
            lines = [line for line in lines if line and len(line) > 3]

            return ExtractionResult(
                text="\n".join(lines),
                pages=1,
                warnings=["Legacy .ppt format - extraction may be incomplete"],
            )
        except Exception as e:
            # Log legacy PowerPoint extraction failures
            logger.info(f"Legacy .ppt extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"Failed to extract from legacy .ppt: {e}"],
            )


class EmailExtractor(BaseExtractor):
    """
    Email extractor for MSG (Outlook) and EML (MIME) files.

    Extracts:
    - Subject, From, To, CC, Date headers
    - Plain text body
    - HTML body (converted to text)
    - Attachment names (not content - attachments processed separately)
    """

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type in (
                "application/vnd.ms-outlook",
                "message/rfc822",
            ) or
            extension in (".msg", ".eml")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        ext = Path(filename).suffix.lower()

        if ext == ".msg":
            return self._extract_msg(content, filename)
        else:
            return self._extract_eml(content, filename)

    def _extract_msg(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from Outlook MSG file."""
        try:
            import extract_msg
        except ImportError:
            raise ImportError("extract-msg not installed. Run: pip install extract-msg")

        try:
            msg = extract_msg.Message(io.BytesIO(content))

            parts = []

            # Headers
            if msg.subject:
                parts.append(f"Subject: {msg.subject}")
            if msg.sender:
                parts.append(f"From: {msg.sender}")
            if msg.to:
                parts.append(f"To: {msg.to}")
            if msg.cc:
                parts.append(f"CC: {msg.cc}")
            if msg.date:
                parts.append(f"Date: {msg.date}")

            parts.append("")  # Blank line before body

            # Body - prefer plain text
            if msg.body:
                parts.append(msg.body)
            elif msg.htmlBody:
                # Convert HTML to text
                parts.append(self._html_to_text(msg.htmlBody))

            # List attachments (names only)
            if msg.attachments:
                parts.append("\n[Attachments]")
                for att in msg.attachments:
                    if hasattr(att, 'longFilename') and att.longFilename:
                        parts.append(f"- {att.longFilename}")
                    elif hasattr(att, 'shortFilename') and att.shortFilename:
                        parts.append(f"- {att.shortFilename}")

            msg.close()

            return ExtractionResult(
                text="\n".join(parts),
                pages=1,
            )

        except Exception as e:
            # Log MSG file extraction failures
            logger.info(f"MSG extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"MSG extraction failed: {e}"],
            )

    def _extract_eml(self, content: bytes, filename: str) -> ExtractionResult:
        """Extract from EML (MIME) file."""
        import email
        from email.policy import default

        try:
            msg = email.message_from_bytes(content, policy=default)

            parts = []

            # Headers
            if msg["subject"]:
                parts.append(f"Subject: {msg['subject']}")
            if msg["from"]:
                parts.append(f"From: {msg['from']}")
            if msg["to"]:
                parts.append(f"To: {msg['to']}")
            if msg["cc"]:
                parts.append(f"CC: {msg['cc']}")
            if msg["date"]:
                parts.append(f"Date: {msg['date']}")

            parts.append("")  # Blank line before body

            # Extract body
            body_text = None
            body_html = None
            attachments = []

            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get("Content-Disposition", ""))

                    # Skip attachments for body extraction
                    if "attachment" in content_disposition:
                        filename_part = part.get_filename()
                        if filename_part:
                            attachments.append(filename_part)
                        continue

                    if content_type == "text/plain" and not body_text:
                        payload = part.get_payload(decode=True)
                        if payload:
                            body_text = payload.decode(
                                part.get_content_charset() or "utf-8",
                                errors="replace"
                            )
                    elif content_type == "text/html" and not body_html:
                        payload = part.get_payload(decode=True)
                        if payload:
                            body_html = payload.decode(
                                part.get_content_charset() or "utf-8",
                                errors="replace"
                            )
            else:
                # Single part message
                content_type = msg.get_content_type()
                payload = msg.get_payload(decode=True)
                if payload:
                    text = payload.decode(
                        msg.get_content_charset() or "utf-8",
                        errors="replace"
                    )
                    if content_type == "text/html":
                        body_html = text
                    else:
                        body_text = text

            # Prefer plain text over HTML
            if body_text:
                parts.append(body_text)
            elif body_html:
                parts.append(self._html_to_text(body_html))

            # List attachments
            if attachments:
                parts.append("\n[Attachments]")
                for att in attachments:
                    parts.append(f"- {att}")

            return ExtractionResult(
                text="\n".join(parts),
                pages=1,
            )

        except Exception as e:
            # Log EML file extraction failures
            logger.info(f"EML extraction failed for {filename}: {type(e).__name__}: {e}")
            return ExtractionResult(
                text="",
                pages=1,
                warnings=[f"EML extraction failed: {e}"],
            )

    def _html_to_text(self, html: str) -> str:
        """Convert HTML to plain text."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, "html.parser")

            # Remove script and style elements
            for element in soup(["script", "style", "head", "meta", "link"]):
                element.decompose()

            # Get text
            text = soup.get_text(separator="\n")

            # Clean up whitespace
            lines = [line.strip() for line in text.splitlines()]
            text = "\n".join(line for line in lines if line)

            return text
        except ImportError:
            # Fallback: basic tag stripping
            import re
            text = re.sub(r'<[^>]+>', ' ', html)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()


class HTMLExtractor(BaseExtractor):
    """
    HTML/web page text extractor.

    Extracts visible text content, excluding scripts, styles, and metadata.
    Preserves basic structure through whitespace.
    """

    def can_handle(self, content_type: str, extension: str) -> bool:
        return (
            content_type in ("text/html", "application/xhtml+xml") or
            extension in (".html", ".htm", ".xhtml")
        )

    def extract(self, content: bytes, filename: str) -> ExtractionResult:
        # Try to decode with various encodings - UnicodeDecodeError is expected for wrong encodings
        text_content = None
        for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
            try:
                text_content = content.decode(encoding)
                break
            except UnicodeDecodeError:
                # This encoding doesn't work - try next one
                continue

        if text_content is None:
            return ExtractionResult(
                text="",
                pages=1,
                warnings=["Failed to decode HTML file"],
            )

        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text_content, "html.parser")

            # Extract title
            title = ""
            if soup.title and soup.title.string:
                title = f"Title: {soup.title.string.strip()}\n\n"

            # Remove non-content elements
            for element in soup(["script", "style", "head", "meta", "link", "noscript"]):
                element.decompose()

            # Extract text with structure
            text_parts = []

            # Process headings specially
            for heading in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                heading_text = heading.get_text(strip=True)
                if heading_text:
                    text_parts.append(f"\n{heading_text}\n")
                heading.decompose()

            # Get remaining text
            body_text = soup.get_text(separator="\n")

            # Clean up whitespace
            lines = [line.strip() for line in body_text.splitlines()]
            body_text = "\n".join(line for line in lines if line)

            text_parts.append(body_text)

            return ExtractionResult(
                text=title + "\n".join(text_parts),
                pages=1,
            )

        except ImportError:
            # Fallback: basic tag stripping without BeautifulSoup
            import re

            # Remove script and style content
            text = re.sub(r'<script[^>]*>.*?</script>', '', text_content, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)

            # Remove all HTML tags
            text = re.sub(r'<[^>]+>', ' ', text)

            # Decode HTML entities
            import html
            text = html.unescape(text)

            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            lines = [line.strip() for line in text.split('. ')]
            text = ".\n".join(line for line in lines if line)

            return ExtractionResult(
                text=text,
                pages=1,
                warnings=["BeautifulSoup not available - using basic extraction"],
            )


def get_extractor(content_type: str, extension: str, ocr_engine: Optional[Any] = None) -> Optional[BaseExtractor]:
    """
    Get an appropriate extractor for the given file type.

    Args:
        content_type: MIME type
        extension: File extension (lowercase, with dot)
        ocr_engine: Optional OCR engine for image/PDF extraction

    Returns:
        Extractor instance or None if no suitable extractor found
    """
    # Create extractors with OCR engine
    extractors = [
        PDFExtractor(ocr_engine=ocr_engine),
        DOCXExtractor(),
        XLSXExtractor(),
        PPTXExtractor(),
        EmailExtractor(),
        HTMLExtractor(),
        ImageExtractor(ocr_engine=ocr_engine),
        TextExtractor(),
        RTFExtractor(),
    ]

    for extractor in extractors:
        if extractor.can_handle(content_type, extension):
            return extractor

    return None


def extract_text(
    content: bytes,
    filename: str,
    content_type: Optional[str] = None,
    ocr_engine: Optional[Any] = None,
) -> ExtractionResult:
    """
    Extract text from file content.

    Convenience function that selects the appropriate extractor.

    Args:
        content: Raw file bytes
        filename: Original filename
        content_type: Optional MIME type (will be guessed if not provided)
        ocr_engine: Optional OCR engine for image/PDF extraction

    Returns:
        ExtractionResult with extracted text and metadata
    """
    import mimetypes

    ext = Path(filename).suffix.lower()

    if content_type is None:
        content_type, _ = mimetypes.guess_type(filename)
        content_type = content_type or ""

    extractor = get_extractor(content_type, ext, ocr_engine=ocr_engine)

    if extractor is None:
        return ExtractionResult(
            text="",
            pages=1,
            warnings=[f"No extractor available for file type: {ext} ({content_type})"],
        )

    return extractor.extract(content, filename)

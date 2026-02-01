"""
File processor for classification pipeline.

Integrates adapters (filesystem, SharePoint, OneDrive) with detection engine
to scan files and produce classification results.

Pipeline:
    1. Fetch file content via adapter
    2. Extract text (based on file type)
    3. Run detection engine
    4. Score entities
    5. Return classification result
"""

import asyncio
import logging
import mimetypes
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Set, AsyncIterator, Union

from .types import Span, DetectionResult, ScoringResult, RiskTier
from .detectors.orchestrator import DetectorOrchestrator
from .scoring.scorer import score
from .constants import DEFAULT_MODELS_DIR

logger = logging.getLogger(__name__)

# Minimum native text length to consider PDF as text-based (not scanned)
MIN_NATIVE_TEXT_LENGTH = 20


# Supported text-based file extensions
TEXT_EXTENSIONS: Set[str] = frozenset({
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv",
    ".json", ".xml", ".yaml", ".yml", ".ini", ".cfg", ".conf",
    ".log", ".sql", ".html", ".htm", ".css", ".js", ".ts",
    ".py", ".rb", ".go", ".java", ".c", ".cpp", ".h", ".hpp",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
})

# Office document extensions (require extraction libraries)
OFFICE_EXTENSIONS: Set[str] = frozenset({
    ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".rtf",
})

# PDF extension
PDF_EXTENSIONS: Set[str] = frozenset({".pdf"})

# Image extensions (require OCR)
IMAGE_EXTENSIONS: Set[str] = frozenset({
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif", ".webp",
})


@dataclass
class FileClassification:
    """Classification result for a single file."""
    file_path: str
    file_name: str
    file_size: int
    mime_type: Optional[str]
    exposure_level: str

    # Detection results
    spans: List[Span] = field(default_factory=list)
    entity_counts: Dict[str, int] = field(default_factory=dict)

    # Scoring results
    risk_score: int = 0
    risk_tier: RiskTier = RiskTier.MINIMAL

    # Metadata
    processed_at: datetime = field(default_factory=lambda: datetime.now(timezone.utc))
    processing_time_ms: float = 0.0
    error: Optional[str] = None

    def to_dict(self) -> dict:
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "file_size": self.file_size,
            "mime_type": self.mime_type,
            "exposure_level": self.exposure_level,
            "entity_counts": self.entity_counts,
            "risk_score": self.risk_score,
            "risk_tier": self.risk_tier.value,
            "processed_at": self.processed_at.isoformat(),
            "processing_time_ms": self.processing_time_ms,
            "error": self.error,
        }


class FileProcessor:
    """
    Processes files through the classification pipeline.

    Usage:
        processor = FileProcessor()

        # Process a single file
        result = await processor.process_file(file_path, content, exposure)

        # Process multiple files
        async for result in processor.process_batch(files):
            print(f"{result.file_path}: {result.risk_tier}")
    """

    def __init__(
        self,
        enable_ml: bool = False,
        enable_ocr: bool = True,
        ml_model_dir: Optional[Path] = None,
        confidence_threshold: float = 0.70,
        max_file_size: int = 50 * 1024 * 1024,  # 50 MB
    ):
        """
        Initialize the processor.

        Args:
            enable_ml: Enable ML-based detectors
            enable_ocr: Enable OCR for images and scanned PDFs
            ml_model_dir: Path to ML model files (defaults to ~/.openlabels/models/)
            confidence_threshold: Minimum detection confidence
            max_file_size: Maximum file size to process (bytes)
        """
        self.max_file_size = max_file_size
        self.enable_ocr = enable_ocr
        self._ocr_engine = None
        self._ml_model_dir = ml_model_dir or DEFAULT_MODELS_DIR
        self._orchestrator = DetectorOrchestrator(
            enable_ml=enable_ml,
            ml_model_dir=ml_model_dir,
            confidence_threshold=confidence_threshold,
        )

        # Lazily initialize OCR engine when needed
        if enable_ocr:
            self._init_ocr_engine()

    def _init_ocr_engine(self) -> None:
        """Initialize OCR engine lazily."""
        try:
            from .ocr import OCREngine
            self._ocr_engine = OCREngine(models_dir=self._ml_model_dir)
            if self._ocr_engine.is_available:
                # Start loading in background for faster first use
                self._ocr_engine.start_loading()
            else:
                logger.info("OCR not available - rapidocr-onnxruntime not installed")
                self._ocr_engine = None
        except ImportError:
            logger.debug("OCR module not available")
            self._ocr_engine = None

    async def process_file(
        self,
        file_path: str,
        content: Union[str, bytes],
        exposure_level: str = "PRIVATE",
        file_size: Optional[int] = None,
    ) -> FileClassification:
        """
        Process a single file.

        Args:
            file_path: Path or identifier for the file
            content: File content (text or bytes)
            exposure_level: File exposure (PRIVATE, INTERNAL, ORG_WIDE, PUBLIC)
            file_size: File size in bytes (for reporting)

        Returns:
            FileClassification with detection and scoring results
        """
        import time
        start_time = time.time()

        file_name = Path(file_path).name
        mime_type, _ = mimetypes.guess_type(file_path)

        # Initialize result
        result = FileClassification(
            file_path=file_path,
            file_name=file_name,
            file_size=file_size or len(content),
            mime_type=mime_type,
            exposure_level=exposure_level,
        )

        try:
            # Extract text if bytes
            if isinstance(content, bytes):
                text = await self._extract_text(content, file_path)
            else:
                text = content

            if not text or not text.strip():
                result.processing_time_ms = (time.time() - start_time) * 1000
                return result

            # Run detection
            detection_result = self._orchestrator.detect(text)
            result.spans = detection_result.spans
            result.entity_counts = detection_result.entity_counts

            # Score entities
            if detection_result.entity_counts:
                score_result = score(
                    entities=detection_result.entity_counts,
                    exposure=exposure_level,
                )
                result.risk_score = score_result.score
                result.risk_tier = score_result.tier

        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            result.error = str(e)

        result.processing_time_ms = (time.time() - start_time) * 1000
        return result

    async def process_batch(
        self,
        files: List[Dict],
        concurrency: int = 4,
    ) -> AsyncIterator[FileClassification]:
        """
        Process multiple files concurrently.

        Args:
            files: List of file dicts with keys: path, content, exposure, size
            concurrency: Maximum concurrent file processing

        Yields:
            FileClassification for each processed file
        """
        semaphore = asyncio.Semaphore(concurrency)

        async def process_one(file_info: Dict) -> FileClassification:
            async with semaphore:
                return await self.process_file(
                    file_path=file_info["path"],
                    content=file_info["content"],
                    exposure_level=file_info.get("exposure", "PRIVATE"),
                    file_size=file_info.get("size"),
                )

        tasks = [process_one(f) for f in files]

        for coro in asyncio.as_completed(tasks):
            result = await coro
            yield result

    async def _extract_text(self, content: bytes, file_path: str) -> str:
        """
        Extract text from file content.

        Args:
            content: Raw file bytes
            file_path: File path for type detection

        Returns:
            Extracted text
        """
        ext = Path(file_path).suffix.lower()

        # Plain text files
        if ext in TEXT_EXTENSIONS:
            return await self._decode_text(content)

        # Office documents
        if ext in OFFICE_EXTENSIONS:
            return await self._extract_office(content, ext)

        # PDFs (with OCR fallback for scanned documents)
        if ext in PDF_EXTENSIONS:
            return await self._extract_pdf(content)

        # Image files (require OCR)
        if ext in IMAGE_EXTENSIONS:
            return await self._extract_image(content)

        # Unknown - try as text
        return await self._decode_text(content)

    async def _extract_image(self, content: bytes) -> str:
        """
        Extract text from image using OCR.

        Args:
            content: Raw image bytes

        Returns:
            Extracted text or empty string if OCR unavailable
        """
        if not self._ocr_engine:
            logger.warning("OCR not available for image extraction")
            return ""

        try:
            import io
            from PIL import Image
            import numpy as np

            # Load image from bytes
            image = Image.open(io.BytesIO(content))
            # Convert to RGB if needed (OCR expects RGB)
            if image.mode != "RGB":
                image = image.convert("RGB")
            # Convert to numpy array for RapidOCR
            image_array = np.array(image)

            # Extract text using OCR
            text = self._ocr_engine.extract_text(image_array)
            return text

        except ImportError as e:
            logger.warning(f"Image processing library not installed: {e}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""

    async def _decode_text(self, content: bytes) -> str:
        """Decode bytes to text with encoding detection."""
        # Try common encodings
        for encoding in ["utf-8", "utf-16", "latin-1", "cp1252"]:
            try:
                return content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                continue

        # Last resort: decode with errors replaced
        return content.decode("utf-8", errors="replace")

    async def _extract_office(self, content: bytes, ext: str) -> str:
        """
        Extract text from Office documents.

        Supports:
        - .docx (python-docx)
        - .xlsx (openpyxl)
        - .pptx (python-pptx)
        - .doc, .xls, .ppt (legacy - limited support)
        - .odt, .ods, .odp (OpenDocument - via zipfile)
        - .rtf (basic extraction)
        """
        import io

        try:
            # Word documents (.docx)
            if ext == ".docx":
                return await self._extract_docx(content)

            # Excel spreadsheets (.xlsx)
            elif ext == ".xlsx":
                return await self._extract_xlsx(content)

            # PowerPoint presentations (.pptx)
            elif ext == ".pptx":
                return await self._extract_pptx(content)

            # OpenDocument formats (.odt, .ods, .odp)
            elif ext in (".odt", ".ods", ".odp"):
                return await self._extract_odf(content)

            # RTF files
            elif ext == ".rtf":
                return await self._extract_rtf(content)

            # Legacy Office formats (.doc, .xls, .ppt)
            elif ext in (".doc", ".xls", ".ppt"):
                logger.warning(f"Legacy Office format {ext} has limited support")
                # Try to extract any embedded text
                return await self._extract_legacy_office(content)

            else:
                logger.warning(f"Unsupported Office format: {ext}")
                return ""

        except ImportError as e:
            logger.warning(f"Office extraction library not installed: {e}")
            return ""
        except Exception as e:
            logger.error(f"Error extracting Office document: {e}")
            return ""

    async def _extract_docx(self, content: bytes) -> str:
        """Extract text from .docx files."""
        try:
            from docx import Document
            import io

            doc = Document(io.BytesIO(content))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-docx not installed. Install with: pip install python-docx")
            # Fallback: try to extract from XML directly
            return await self._extract_docx_fallback(content)

    async def _extract_docx_fallback(self, content: bytes) -> str:
        """Fallback extraction from .docx using zipfile."""
        import zipfile
        import io
        import re

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "word/document.xml" in zf.namelist():
                    xml_content = zf.read("word/document.xml").decode("utf-8")
                    # Strip XML tags
                    text = re.sub(r"<[^>]+>", " ", xml_content)
                    # Clean up whitespace
                    text = re.sub(r"\s+", " ", text).strip()
                    return text
        except Exception as e:
            logger.debug(f"Fallback docx extraction failed: {e}")
        return ""

    async def _extract_xlsx(self, content: bytes) -> str:
        """Extract text from .xlsx files."""
        try:
            from openpyxl import load_workbook
            import io

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []

            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text_parts.append(" ".join(row_text))

            wb.close()
            return "\n".join(text_parts)

        except ImportError:
            logger.warning("openpyxl not installed. Install with: pip install openpyxl")
            # Fallback: try to extract from XML directly
            return await self._extract_xlsx_fallback(content)

    async def _extract_xlsx_fallback(self, content: bytes) -> str:
        """Fallback extraction from .xlsx using zipfile."""
        import zipfile
        import io
        import re

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                text_parts = []
                for name in zf.namelist():
                    if name.startswith("xl/worksheets/") and name.endswith(".xml"):
                        xml_content = zf.read(name).decode("utf-8")
                        # Extract values from <v> tags
                        values = re.findall(r"<v>([^<]+)</v>", xml_content)
                        text_parts.extend(values)
                return " ".join(text_parts)
        except Exception as e:
            logger.debug(f"Fallback xlsx extraction failed: {e}")
        return ""

    async def _extract_pptx(self, content: bytes) -> str:
        """Extract text from .pptx files."""
        try:
            from pptx import Presentation
            import io

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            # Fallback: try to extract from XML directly
            return await self._extract_pptx_fallback(content)

    async def _extract_pptx_fallback(self, content: bytes) -> str:
        """Fallback extraction from .pptx using zipfile."""
        import zipfile
        import io
        import re

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                text_parts = []
                for name in zf.namelist():
                    if name.startswith("ppt/slides/") and name.endswith(".xml"):
                        xml_content = zf.read(name).decode("utf-8")
                        # Extract text from <a:t> tags
                        texts = re.findall(r"<a:t>([^<]+)</a:t>", xml_content)
                        text_parts.extend(texts)
                return " ".join(text_parts)
        except Exception as e:
            logger.debug(f"Fallback pptx extraction failed: {e}")
        return ""

    async def _extract_odf(self, content: bytes) -> str:
        """Extract text from OpenDocument formats (.odt, .ods, .odp)."""
        import zipfile
        import io
        import re

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "content.xml" in zf.namelist():
                    xml_content = zf.read("content.xml").decode("utf-8")
                    # Strip XML tags
                    text = re.sub(r"<[^>]+>", " ", xml_content)
                    # Clean up whitespace
                    text = re.sub(r"\s+", " ", text).strip()
                    return text
        except Exception as e:
            logger.debug(f"ODF extraction failed: {e}")
        return ""

    async def _extract_rtf(self, content: bytes) -> str:
        """Extract text from RTF files."""
        import re

        try:
            # Decode RTF content
            text = content.decode("latin-1")

            # Remove RTF control words and groups
            text = re.sub(r"\\[a-z]+\d*\s?", " ", text)
            text = re.sub(r"[{}]", "", text)
            text = re.sub(r"\\[^a-z]", "", text)

            # Clean up whitespace
            text = re.sub(r"\s+", " ", text).strip()

            return text
        except Exception as e:
            logger.debug(f"RTF extraction failed: {e}")
        return ""

    async def _extract_legacy_office(self, content: bytes) -> str:
        """
        Extract text from legacy Office formats (.doc, .xls, .ppt).

        These formats are binary (OLE2) and require specialized libraries.
        This provides basic extraction for any plaintext embedded in the file.
        """
        import re

        try:
            # Try to decode as text (may contain readable strings)
            text = content.decode("latin-1", errors="ignore")

            # Extract printable ASCII sequences (4+ chars)
            strings = re.findall(r"[\x20-\x7e]{4,}", text)

            # Filter out obvious binary garbage
            filtered = [s for s in strings if not re.match(r"^[\x00-\x1f]+$", s)]

            return " ".join(filtered[:1000])  # Limit to prevent huge outputs
        except Exception as e:
            logger.debug(f"Legacy Office extraction failed: {e}")
        return ""

    async def _extract_pdf(self, content: bytes) -> str:
        """
        Extract text from PDF files.

        Tries multiple extraction methods in order:
        1. pdfplumber (best for structured PDFs)
        2. PyMuPDF/fitz (fast, good for scanned PDFs with OCR)
        3. pypdf (lightweight fallback)
        4. OCR fallback for scanned documents (if native text < MIN_NATIVE_TEXT_LENGTH)
        """
        import io

        native_text = ""

        # Try pdfplumber first (best quality)
        try:
            import pdfplumber

            text_parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            if text_parts:
                native_text = "\n".join(text_parts)

        except ImportError:
            logger.debug("pdfplumber not installed, trying alternatives")
        except Exception as e:
            logger.debug(f"pdfplumber extraction failed: {e}")

        # Try PyMuPDF (fitz) second
        if not native_text:
            try:
                import fitz  # PyMuPDF

                text_parts = []
                doc = fitz.open(stream=content, filetype="pdf")

                for page in doc:
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(text)

                doc.close()

                if text_parts:
                    native_text = "\n".join(text_parts)

            except ImportError:
                logger.debug("PyMuPDF not installed, trying alternatives")
            except Exception as e:
                logger.debug(f"PyMuPDF extraction failed: {e}")

        # Try pypdf as fallback
        if not native_text:
            try:
                from pypdf import PdfReader

                text_parts = []
                reader = PdfReader(io.BytesIO(content))

                for page in reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text)

                if text_parts:
                    native_text = "\n".join(text_parts)

            except ImportError:
                logger.debug("pypdf not installed, trying alternatives")
            except Exception as e:
                logger.debug(f"pypdf extraction failed: {e}")

        # Last resort: try PyPDF2 (older but common)
        if not native_text:
            try:
                from PyPDF2 import PdfReader as PyPDF2Reader

                text_parts = []
                reader = PyPDF2Reader(io.BytesIO(content))

                for page in reader.pages:
                    text = page.extract_text()
                    if text and text.strip():
                        text_parts.append(text)

                if text_parts:
                    native_text = "\n".join(text_parts)

            except ImportError:
                logger.warning(
                    "No PDF library installed. Install one of: "
                    "pip install pdfplumber OR pip install pymupdf OR pip install pypdf"
                )
            except Exception as e:
                logger.debug(f"PyPDF2 extraction failed: {e}")

        # If native text is too short, try OCR (likely a scanned document)
        if len(native_text.strip()) < MIN_NATIVE_TEXT_LENGTH and self._ocr_engine:
            logger.info("PDF has minimal native text, attempting OCR extraction")
            ocr_text = await self._extract_pdf_with_ocr(content)
            if ocr_text:
                return ocr_text

        return native_text

    async def _extract_pdf_with_ocr(self, content: bytes) -> str:
        """
        Extract text from scanned PDF pages using OCR.

        Args:
            content: PDF file bytes

        Returns:
            OCR-extracted text from all pages
        """
        if not self._ocr_engine:
            return ""

        try:
            import io

            # Try PyMuPDF for rendering pages to images
            try:
                import fitz
                import numpy as np

                text_parts = []
                doc = fitz.open(stream=content, filetype="pdf")

                for page_num, page in enumerate(doc):
                    # Render page to image (150 DPI for good OCR quality)
                    mat = fitz.Matrix(150 / 72, 150 / 72)  # 150 DPI
                    pix = page.get_pixmap(matrix=mat)

                    # Convert to numpy array
                    img_array = np.frombuffer(pix.samples, dtype=np.uint8)
                    img_array = img_array.reshape(pix.height, pix.width, pix.n)

                    # Convert to RGB if RGBA
                    if pix.n == 4:
                        img_array = img_array[:, :, :3]

                    # Run OCR
                    page_text = self._ocr_engine.extract_text(img_array)
                    if page_text:
                        text_parts.append(page_text)

                doc.close()

                if text_parts:
                    return "\n\n".join(text_parts)

            except ImportError:
                logger.debug("PyMuPDF not available for PDF OCR rendering")

            # Fallback: try pdf2image if available
            try:
                from pdf2image import convert_from_bytes
                import numpy as np

                images = convert_from_bytes(content, dpi=150)
                text_parts = []

                for img in images:
                    img_array = np.array(img)
                    page_text = self._ocr_engine.extract_text(img_array)
                    if page_text:
                        text_parts.append(page_text)

                if text_parts:
                    return "\n\n".join(text_parts)

            except ImportError:
                logger.debug("pdf2image not available for PDF OCR rendering")

        except Exception as e:
            logger.error(f"Error extracting PDF with OCR: {e}")

        return ""

    def can_process(self, file_path: str, file_size: int) -> bool:
        """
        Check if a file can be processed.

        Args:
            file_path: Path to the file
            file_size: Size in bytes

        Returns:
            True if file can be processed
        """
        if file_size > self.max_file_size:
            return False

        ext = Path(file_path).suffix.lower()
        supported = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | PDF_EXTENSIONS | IMAGE_EXTENSIONS

        return ext in supported


async def process_file(
    file_path: str,
    content: Union[str, bytes],
    exposure_level: str = "PRIVATE",
    **kwargs,
) -> FileClassification:
    """
    Convenience function to process a single file.

    Args:
        file_path: Path to the file
        content: File content
        exposure_level: File exposure level
        **kwargs: Passed to FileProcessor

    Returns:
        FileClassification result
    """
    processor = FileProcessor(**kwargs)
    return await processor.process_file(file_path, content, exposure_level)
